from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colors ──────────────────────────────────────────
PRIMARY      = RGBColor(0x8B, 0x18, 0x32)
PRIMARY_DARK = RGBColor(0x5C, 0x0F, 0x21)
PRIMARY_LITE = RGBColor(0xB2, 0x2C, 0x4A)
GOLD         = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_LITE    = RGBColor(0xE8, 0xC9, 0x7A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF8, 0xF5, 0xF2)
GRAY_LIGHT   = RGBColor(0xF5, 0xF5, 0xF7)
GRAY         = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY    = RGBColor(0x66, 0x66, 0x66)
TEXT_DARK    = RGBColor(0x1A, 0x1A, 0x1A)
GREEN        = RGBColor(0x2E, 0x7D, 0x32)

# ── Slide size: 16:9 widescreen ─────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # blank layout

# ═══════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def add_textbox_ml(slide, lines, x, y, w, h):
    """Multi-line textbox. lines = list of (text, size, bold, color, align)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txb

def slide_number(slide, num, total, light=True):
    color = WHITE if light else DARK_GRAY
    add_text(slide, f"{num:02d} / {total}", W - Inches(1.5), H - Inches(0.4),
             Inches(1.3), Inches(0.35), font_size=10, color=color,
             align=PP_ALIGN.RIGHT)
    wm_color = RGBColor(0xCC, 0xCC, 0xCC) if not light else RGBColor(0xAA, 0xAA, 0xAA)
    add_text(slide, "ATB MOBILE · 2026", Inches(0.2), H - Inches(0.4),
             Inches(2.5), Inches(0.35), font_size=9, bold=True,
             color=wm_color if light else DARK_GRAY, align=PP_ALIGN.LEFT)

def add_card(slide, x, y, w, h, bg=None):
    bg = bg or GRAY_LIGHT
    shape = add_rect(slide, x, y, w, h, bg)
    # rounded effect via border
    sp = shape._element
    return shape

def eyebrow(slide, text, x, y, w=Inches(8), color=PRIMARY):
    add_text(slide, text.upper(), x, y, w, Inches(0.3),
             font_size=10, bold=True, color=color, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# background gradient simulation — two rects
add_rect(s1, 0, 0, W, H, PRIMARY_DARK)
add_rect(s1, 0, 0, W * 0.65, H, PRIMARY)

# gold accent bar top
add_rect(s1, 0, 0, W, Inches(0.06), GOLD)

# right decorative circle
circ = s1.shapes.add_shape(9, W - Inches(4.5), -Inches(1.5), Inches(5), Inches(5))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
circ.line.fill.background()
circ._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
sp_pr = circ._element.spPr
a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
solid = sp_pr.find(f'{{{a_ns}}}solidFill')
if solid is not None:
    srgb = solid.find(f'{{{a_ns}}}srgbClr')
    if srgb is not None:
        alpha_el = etree.SubElement(srgb, f'{{{a_ns}}}alpha')
        alpha_el.set('val', '8000')  # ~5% opacity

# ATB badge box
add_rect(s1, Inches(0.6), Inches(0.8), Inches(2.2), Inches(0.45),
         RGBColor(0xFF, 0xFF, 0xFF))
add_text(s1, "ATB", Inches(0.65), Inches(0.82), Inches(0.55), Inches(0.38),
         font_size=11, bold=True, color=PRIMARY)
add_text(s1, "Arab Tunisian Bank", Inches(1.25), Inches(0.82), Inches(1.6), Inches(0.38),
         font_size=10, bold=False, color=PRIMARY_DARK)

# Main title
add_textbox_ml(s1, [
    ("Banque Mobile", 52, True, WHITE, PP_ALIGN.LEFT),
    ("Nouvelle Ère", 52, True, GOLD_LITE, PP_ALIGN.LEFT),
], Inches(0.6), Inches(1.5), Inches(6.5), Inches(2.2))

# Subtitle
add_text(s1, "Une application bancaire complète pour les clients ATB —\nDABs en temps réel, assistant IA en Derja,\net gestion de cartes intuitive.",
         Inches(0.6), Inches(3.6), Inches(6.0), Inches(1.2),
         font_size=15, color=RGBColor(0xCC, 0xCC, 0xCC))

# Pills row
pills = ["📍  Carte des DABs", "🤖  Assistant IA", "💳  Cartes", "🔔  Notifications"]
px = Inches(0.6)
for pill in pills:
    add_rect(s1, px, Inches(5.1), Inches(2.05), Inches(0.42), RGBColor(0x6A, 0x10, 0x22))
    add_text(s1, pill, px + Inches(0.12), Inches(5.12), Inches(1.85), Inches(0.4),
             font_size=11, color=WHITE)
    px += Inches(2.15)

# Phone mockup (simplified)
add_rect(s1, Inches(9.2), Inches(1.0), Inches(1.8), Inches(5.2), RGBColor(0x1A, 0x1A, 0x1A))
add_rect(s1, Inches(9.3), Inches(1.5), Inches(1.6), Inches(4.2), PRIMARY_DARK)
# screen elements
for iy, ih in [(1.7, 0.12), (1.95, 0.08), (2.15, 0.08), (2.5, 0.7), (3.3, 0.08), (3.5, 0.08)]:
    add_rect(s1, Inches(9.45), Inches(iy), Inches(1.3), Inches(ih),
             RGBColor(0xFF, 0xFF, 0xFF))

add_rect(s1, Inches(11.2), Inches(1.3), Inches(1.8), Inches(4.6), RGBColor(0x1A, 0x1A, 0x1A))
add_rect(s1, Inches(11.3), Inches(1.75), Inches(1.6), Inches(3.7), RGBColor(0x1A, 0x3A, 0x1A))

slide_number(s1, 1, 10, light=True)

# ═══════════════════════════════════════════════════
#  SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, W, H, WHITE)

# Header bar
add_rect(s2, 0, 0, W, Inches(1.1), PRIMARY_DARK)
add_rect(s2, 0, Inches(1.08), W, Inches(0.04), GOLD)
add_text(s2, "📊  Résumé Exécutif", Inches(0.5), Inches(0.3),
         Inches(6), Inches(0.6), font_size=22, bold=True, color=WHITE)
add_text(s2, "ATB", W - Inches(1.4), Inches(0.3), Inches(1.2), Inches(0.55),
         font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left column title
eyebrow(s2, "Vision du Projet", Inches(0.5), Inches(1.3), color=PRIMARY)

points = [
    ("1", "Application 100% Flutter — une seule base de code Android & iOS,\nréduisant les coûts de développement de 40%."),
    ("2", "Cartographie temps réel des DABs ATB via OpenStreetMap —\naucune dépendance à une API payante."),
    ("3", "Assistant IA multilingue (Darija · Français · Anglais)\npropulsé par Groq LLaMA, réponses en moins de 500ms."),
    ("4", "Gestion complète des cartes bancaires avec alertes sécurité,\nlimites et historique des transactions."),
    ("5", "Architecture modulaire prête pour l'intégration\navec les APIs backend ATB."),
]

py = Inches(1.65)
for num, text in points:
    add_rect(s2, Inches(0.5), py + Inches(0.02), Inches(0.36), Inches(0.36), PRIMARY)
    add_text(s2, num, Inches(0.5), py, Inches(0.36), Inches(0.4),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, text, Inches(0.97), py - Inches(0.02), Inches(5.6), Inches(0.52),
             font_size=12, color=TEXT_DARK)
    py += Inches(0.72)

# Right column
add_rect(s2, Inches(7.1), Inches(1.2), Inches(5.7), Inches(5.9), GRAY_LIGHT)
eyebrow(s2, "Chiffres Clés", Inches(7.3), Inches(1.3), w=Inches(4), color=PRIMARY)

kpis = [
    ("📱", "5", "Écrans principaux"),
    ("📍", "20+", "DABs pré-chargés"),
    ("⚡", "<500ms", "Temps réponse IA"),
    ("🌍", "3", "Langues supportées"),
]
ky = Inches(1.65)
for icon, val, lbl in kpis:
    add_rect(s2, Inches(7.3), ky, Inches(5.3), Inches(0.75), WHITE)
    add_text(s2, icon, Inches(7.4), ky + Inches(0.1), Inches(0.5), Inches(0.55), font_size=20, color=PRIMARY)
    add_text(s2, val, Inches(8.0), ky + Inches(0.08), Inches(1.8), Inches(0.4),
             font_size=18, bold=True, color=PRIMARY)
    add_text(s2, lbl, Inches(9.65), ky + Inches(0.18), Inches(2.8), Inches(0.35),
             font_size=12, color=DARK_GRAY)
    ky += Inches(0.88)

# Quote box
add_rect(s2, Inches(7.3), Inches(5.3), Inches(5.3), Inches(1.4), PRIMARY)
add_rect(s2, Inches(7.3), Inches(5.3), Inches(0.06), Inches(1.4), GOLD)
add_text(s2, '"Une expérience bancaire moderne, pensée pour\nle client tunisien — rapide, intelligente et accessible."',
         Inches(7.5), Inches(5.45), Inches(4.9), Inches(1.1),
         font_size=12, italic=True, color=WHITE)

slide_number(s2, 2, 10, light=False)

# ═══════════════════════════════════════════════════
#  SLIDE 3 — FEATURE GRID
# ═══════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, W, H, OFF_WHITE)

eyebrow(s3, "Fonctionnalités", Inches(0.6), Inches(0.35), color=PRIMARY)
add_text(s3, "Ce que fait l'application", Inches(0.6), Inches(0.6),
         Inches(7), Inches(0.7), font_size=30, bold=True, color=TEXT_DARK)

features = [
    (PRIMARY, WHITE, "📍", "Carte des DABs en direct",
     "Localisation GPS en temps réel. Markers colorés ouvert/fermé, filtres dépôt et retrait sans carte."),
    (WHITE, TEXT_DARK, "🤖", "Assistant IA · Derja",
     "Chatbot LLaMA 3.1 via Groq. Répond en tunisien, français ou anglais selon l'utilisateur."),
    (WHITE, TEXT_DARK, "💳", "Gestion des Cartes",
     "Carousel VISA/Mastercard, sécurité, limites de retrait et historique complet."),
    (WHITE, TEXT_DARK, "🔔", "Messages & Alertes",
     "Centre de notifications bancaires groupées par date. Filtres Banque · Sécurité · Contacts."),
    (PRIMARY_DARK, WHITE, "🛡️", "Sécurité Avancée",
     "Blocage de carte instantané, alertes de connexion suspecte et authentification biométrique."),
    (WHITE, TEXT_DARK, "🗺️", "Itinéraire Intégré",
     "Ouverture directe dans Google Maps depuis la fiche DAB. Distance calculée en temps réel."),
]

cols = 3
card_w = Inches(3.9)
card_h = Inches(1.9)
gx, gy = Inches(0.55), Inches(1.55)
gap_x, gap_y = Inches(0.22), Inches(0.18)

for i, (bg, tc, icon, title, desc) in enumerate(features):
    cx = gx + (i % cols) * (card_w + gap_x)
    cy = gy + (i // cols) * (card_h + gap_y)
    add_rect(s3, cx, cy, card_w, card_h, bg)
    add_text(s3, icon, cx + Inches(0.15), cy + Inches(0.12), Inches(0.5), Inches(0.45), font_size=22, color=tc)
    add_text(s3, title, cx + Inches(0.7), cy + Inches(0.15), card_w - Inches(0.85), Inches(0.42),
             font_size=13, bold=True, color=tc)
    desc_color = RGBColor(0xCC, 0xCC, 0xCC) if bg in (PRIMARY, PRIMARY_DARK) else DARK_GRAY
    add_text(s3, desc, cx + Inches(0.15), cy + Inches(0.65), card_w - Inches(0.25), Inches(1.1),
             font_size=11, color=desc_color, wrap=True)

slide_number(s3, 3, 10, light=False)

# ═══════════════════════════════════════════════════
#  SLIDE 4 — DAB LOCATOR
# ═══════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, W, H, WHITE)

# Left panel
add_rect(s4, 0, 0, Inches(4.8), H, PRIMARY_DARK)
add_rect(s4, 0, 0, Inches(4.8), H, PRIMARY)
add_rect(s4, 0, 0, Inches(0.06), H, GOLD)

add_text(s4, "FONCTIONNALITÉ PHARE", Inches(0.35), Inches(0.6),
         Inches(4.2), Inches(0.35), font_size=10, bold=True, color=GOLD_LITE)
add_text(s4, "Carte des\nDABs ATB", Inches(0.35), Inches(1.0),
         Inches(4.2), Inches(1.5), font_size=34, bold=True, color=WHITE)
add_text(s4, "Trouvez instantanément le distributeur le plus\nproche, avec filtres en temps réel et navigation intégrée.",
         Inches(0.35), Inches(2.55), Inches(4.2), Inches(0.85),
         font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))

features_list = [
    "✅  Données OpenStreetMap (gratuit)",
    "✅  Filtres : Ouvert · Dépôt · Sans carte",
    "✅  Marker sélectionné avec carte infos",
    "✅  Itinéraire Google Maps en 1 clic",
    "✅  Position GPS de l'utilisateur",
    "✅  20+ DABs pré-chargés (fallback)",
]
fy = Inches(3.55)
for f in features_list:
    add_text(s4, f, Inches(0.35), fy, Inches(4.3), Inches(0.38),
             font_size=12, color=WHITE)
    fy += Inches(0.45)

# Right panel — Map mockup
add_rect(s4, Inches(5.0), Inches(0.35), Inches(7.9), Inches(5.5), RGBColor(0xD4, 0xED, 0xDA))

# Grid lines (roads)
for iy in [2.0, 3.5]:
    add_rect(s4, Inches(5.0), Inches(iy), Inches(7.9), Inches(0.12), WHITE)
for ix in [7.5, 9.8]:
    add_rect(s4, Inches(ix), Inches(0.35), Inches(0.12), Inches(5.5), WHITE)

# ATM pins
pins = [
    (Inches(7.0), Inches(2.3), "ATB Marine", PRIMARY),
    (Inches(9.3), Inches(1.5), "ATB Lac", GREEN),
    (Inches(8.2), Inches(3.8), "BIAT", DARK_GRAY),
    (Inches(10.5), Inches(2.8), "STB", GREEN),
]
for px2, py2, name, color in pins:
    add_rect(s4, px2 - Inches(0.2), py2 - Inches(0.22), Inches(0.4), Inches(0.4), color)
    add_rect(s4, px2 - Inches(0.55), py2 + Inches(0.22), Inches(1.1), Inches(0.28), WHITE)
    add_text(s4, name, px2 - Inches(0.5), py2 + Inches(0.23), Inches(1.0), Inches(0.26),
             font_size=8, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# User dot
add_rect(s4, Inches(7.2), Inches(2.5), Inches(0.22), Inches(0.22), RGBColor(0x21, 0x96, 0xF3))

# ATM info card at bottom
add_rect(s4, Inches(5.0), Inches(6.0), Inches(7.9), Inches(1.15), WHITE)
add_rect(s4, Inches(5.0), Inches(6.0), Inches(7.9), Inches(0.04), PRIMARY)
add_rect(s4, Inches(5.15), Inches(6.15), Inches(0.6), Inches(0.6), GRAY_LIGHT)
add_text(s4, "📍", Inches(5.15), Inches(6.15), Inches(0.6), Inches(0.6),
         font_size=20, color=PRIMARY, align=PP_ALIGN.CENTER)
add_text(s4, "Agence Tunis Marine", Inches(5.85), Inches(6.18),
         Inches(4), Inches(0.35), font_size=14, bold=True, color=TEXT_DARK)
add_text(s4, "Avenue Habib Bourguiba, Tunis  ·  Dépôt ✓  ·  24/7",
         Inches(5.85), Inches(6.55), Inches(5), Inches(0.3),
         font_size=11, color=DARK_GRAY)
add_text(s4, "320m", Inches(11.5), Inches(6.25), Inches(1.2), Inches(0.4),
         font_size=16, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT)

slide_number(s4, 4, 10, light=False)

# ═══════════════════════════════════════════════════
#  SLIDE 5 — AI CHATBOT
# ═══════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, W, H, OFF_WHITE)

# Left
eyebrow(s5, "Intelligence Artificielle", Inches(0.5), Inches(0.4), color=PRIMARY)
add_text(s5, "Assistant Tunisien\nPropulsé par IA", Inches(0.5), Inches(0.7),
         Inches(6.5), Inches(1.6), font_size=34, bold=True, color=TEXT_DARK)
add_text(s5, "Le premier chatbot bancaire tunisien qui comprend la Darija,\nle français et l'anglais — et mélange les trois naturellement.",
         Inches(0.5), Inches(2.35), Inches(6.5), Inches(0.85),
         font_size=13, color=DARK_GRAY)

# Language pills
langs = [("🇹🇳  Derja Tounsia", PRIMARY), ("🇫🇷  Français", RGBColor(0x15, 0x65, 0xC0)), ("🇬🇧  English", GREEN)]
lx = Inches(0.5)
for lbl, clr in langs:
    add_rect(s5, lx, Inches(3.3), Inches(1.85), Inches(0.4), WHITE)
    add_text(s5, lbl, lx + Inches(0.1), Inches(3.32), Inches(1.65), Inches(0.36),
             font_size=11, bold=True, color=clr)
    lx += Inches(2.0)

caps = [
    "📍  Localisation du DAB le plus proche",
    "🕐  Vérification des horaires d'ouverture",
    "💳  Aide au blocage de carte",
    "💰  Info commissions & limites",
    "⚡  Réponses en moins de 500ms (Groq)",
]
cy2 = Inches(3.95)
for c in caps:
    add_text(s5, c, Inches(0.5), cy2, Inches(6.3), Inches(0.38),
             font_size=13, color=TEXT_DARK)
    cy2 += Inches(0.48)

# Right — Chat UI
add_rect(s5, Inches(7.0), Inches(0.2), Inches(5.9), Inches(7.1), WHITE)
# Chat header
add_rect(s5, Inches(7.0), Inches(0.2), Inches(5.9), Inches(0.95), PRIMARY_DARK)
add_rect(s5, Inches(7.0), Inches(0.2), Inches(0.08), Inches(0.95), GOLD)
add_rect(s5, Inches(7.15), Inches(0.32), Inches(0.62), Inches(0.62), RGBColor(0x6A, 0x10, 0x22))
add_text(s5, "🤖", Inches(7.15), Inches(0.28), Inches(0.62), Inches(0.62),
         font_size=20, align=PP_ALIGN.CENTER, color=WHITE)
add_text(s5, "Assistance IA · ATB", Inches(7.85), Inches(0.28),
         Inches(4.8), Inches(0.38), font_size=14, bold=True, color=WHITE)
add_text(s5, "🟢  En ligne pour vous aider", Inches(7.85), Inches(0.62),
         Inches(4.8), Inches(0.3), font_size=10, color=RGBColor(0xCC, 0xCC, 0xCC))

# Chat messages
msgs = [
    ("bot",  "Aslema! Ana assistant ATB. Kifesh naaounek lyoum? 😊"),
    ("user", "Feen el DAB el akreb m3a service dépôt?"),
    ("bot",  "Barra! \"Agence Tunis Marine\" 320m men andek 3andha\nmodule dépôt. Maftouha 08h-18h. 📍 Tħeb l'itinéraire?"),
    ("user", "Ouiii, merci barra!"),
    ("bot",  "Afw! Baraka Lahu fik 🙏 Navigation lancée vers l'agence."),
]
my = Inches(1.35)
for role, text in msgs:
    is_bot = role == "bot"
    bg_clr = GRAY_LIGHT if is_bot else PRIMARY
    tc2 = TEXT_DARK if is_bot else WHITE
    mx = Inches(7.15) if is_bot else Inches(9.7)
    bw = Inches(3.5)
    lines = text.count('\n') + 1
    bh = Inches(0.38 + 0.28 * (lines - 1))
    add_rect(s5, mx, my, bw, bh, bg_clr)
    add_text(s5, text, mx + Inches(0.1), my + Inches(0.06), bw - Inches(0.18), bh - Inches(0.1),
             font_size=11, color=tc2, wrap=True)
    my += bh + Inches(0.2)

# Input bar
add_rect(s5, Inches(7.0), Inches(6.65), Inches(5.9), Inches(0.65), WHITE)
add_rect(s5, Inches(7.0), Inches(6.65), Inches(5.9), Inches(0.04), GRAY_LIGHT)
add_rect(s5, Inches(7.1), Inches(6.73), Inches(4.5), Inches(0.45), GRAY_LIGHT)
add_text(s5, "Posez votre question...", Inches(7.2), Inches(6.77),
         Inches(4.2), Inches(0.35), font_size=11, color=GRAY)
add_rect(s5, Inches(11.72), Inches(6.73), Inches(0.45), Inches(0.45), PRIMARY)
add_text(s5, "▶", Inches(11.72), Inches(6.73), Inches(0.45), Inches(0.45),
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(s5, 5, 10, light=False)

# ═══════════════════════════════════════════════════
#  SLIDE 6 — CARDS
# ═══════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, W, H, WHITE)
add_rect(s6, 0, 0, W, Inches(0.06), GOLD)

eyebrow(s6, "Gestion Bancaire", Inches(0.5), Inches(0.25), color=PRIMARY)
add_text(s6, "Mes Cartes", Inches(0.5), Inches(0.55),
         Inches(7), Inches(0.7), font_size=32, bold=True, color=TEXT_DARK)

# Bank card stack
add_rect(s6, Inches(0.7), Inches(1.6), Inches(3.8), Inches(2.2), PRIMARY_DARK)
add_text(s6, "VISA", Inches(3.6), Inches(1.72), Inches(0.7), Inches(0.45),
         font_size=16, bold=True, color=WHITE)
add_text(s6, "**** **** **** 7891", Inches(0.9), Inches(2.6),
         Inches(3.2), Inches(0.45), font_size=14, color=WHITE)

add_rect(s6, Inches(0.4), Inches(1.3), Inches(3.8), Inches(2.2), PRIMARY)
add_rect(s6, Inches(0.4), Inches(1.3), Inches(0.7), Inches(0.45), GOLD)
add_text(s6, "VISA", Inches(3.3), Inches(1.42), Inches(0.7), Inches(0.45),
         font_size=16, bold=True, color=WHITE)
add_text(s6, "**** **** **** 4242", Inches(0.6), Inches(2.3),
         Inches(3.2), Inches(0.45), font_size=14, color=WHITE)
add_text(s6, "SAMI BEN ALI", Inches(0.6), Inches(3.0),
         Inches(2.0), Inches(0.3), font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(s6, "12/26", Inches(2.8), Inches(3.0),
         Inches(1.2), Inches(0.3), font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC))

# Services grid
svcs = [
    ("📱", "Paiement Mobile",    "Paiements NFC et QR code depuis votre smartphone."),
    ("🛡️", "Sécurité Cartes",    "Blocage instantané, alertes fraude et 2FA."),
    ("🕐", "Historique",         "Toutes vos transactions classées et filtrables."),
    ("⚙️", "Limites",            "Contrôle des plafonds de retrait et de paiement."),
]
sx0, sy0 = Inches(4.8), Inches(1.2)
sw, sh = Inches(4.0), Inches(1.45)
for i, (ic, tt, dd) in enumerate(svcs):
    sx = sx0 + (i % 2) * (sw + Inches(0.2))
    sy = sy0 + (i // 2) * (sh + Inches(0.15))
    add_rect(s6, sx, sy, sw, sh, OFF_WHITE)
    add_text(s6, ic, sx + Inches(0.15), sy + Inches(0.15), Inches(0.5), Inches(0.5),
             font_size=24, color=PRIMARY)
    add_text(s6, tt, sx + Inches(0.75), sy + Inches(0.18), sw - Inches(0.85), Inches(0.38),
             font_size=13, bold=True, color=TEXT_DARK)
    add_text(s6, dd, sx + Inches(0.15), sy + Inches(0.65), sw - Inches(0.25), Inches(0.65),
             font_size=11, color=DARK_GRAY, wrap=True)

slide_number(s6, 6, 10, light=False)

# ═══════════════════════════════════════════════════
#  SLIDE 7 — TECH STACK
# ═══════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, W, H, PRIMARY_DARK)
add_rect(s7, 0, 0, W, Inches(0.06), GOLD)

add_text(s7, "ARCHITECTURE TECHNIQUE", Inches(0.6), Inches(0.3),
         Inches(8), Inches(0.35), font_size=10, bold=True, color=GOLD_LITE)
add_text(s7, "Stack Technologique", Inches(0.6), Inches(0.6),
         Inches(8), Inches(0.7), font_size=32, bold=True, color=WHITE)

techs = [
    ("🐦", "Flutter 3.44",      "Framework UI cross-platform.\nUn code, Android + iOS.",       "Frontend"),
    ("🗺️", "Flutter Map",       "Tuiles OpenStreetMap.\nAucune API key. Gratuit illimité.",     "Cartographie"),
    ("⚡", "Groq + LLaMA 3.1", "Inférence ultra-rapide.\nRéponses en <500ms. API gratuite.",  "IA · Gratuit"),
    ("📡", "Overpass API",      "Récupération des DABs réels\ndepuis OpenStreetMap.",           "Données"),
]

tw = Inches(2.95)
th = Inches(2.8)
tx0 = Inches(0.55)
ty0 = Inches(1.5)
for i, (ic, nm, detail, badge) in enumerate(techs):
    tx = tx0 + i * (tw + Inches(0.2))
    add_rect(s7, tx, ty0, tw, th, RGBColor(0x40, 0x08, 0x15))
    add_rect(s7, tx, ty0, tw, Inches(0.06), GOLD)
    add_text(s7, ic, tx + Inches(0.2), ty0 + Inches(0.2), Inches(0.6), Inches(0.6),
             font_size=28, color=WHITE)
    add_text(s7, nm, tx + Inches(0.2), ty0 + Inches(0.9), tw - Inches(0.3), Inches(0.45),
             font_size=15, bold=True, color=WHITE)
    add_text(s7, detail, tx + Inches(0.2), ty0 + Inches(1.4), tw - Inches(0.3), Inches(0.9),
             font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC), wrap=True)
    add_rect(s7, tx + Inches(0.2), ty0 + Inches(2.35), Inches(1.5), Inches(0.32), RGBColor(0x6A, 0x10, 0x22))
    add_text(s7, badge, tx + Inches(0.25), ty0 + Inches(2.36), Inches(1.4), Inches(0.28),
             font_size=9, bold=True, color=GOLD_LITE)

# Architecture flow
arch_items = ["📱 App Flutter", "→", "🌐 HTTP / REST", "→", "🤖 Groq API",
              "+", "🗺️ Overpass", "+", "📍 GPS Device"]
ax = Inches(0.55)
ay = Inches(4.6)
for item in arch_items:
    is_arr = item in ("→", "+")
    if not is_arr:
        add_rect(s7, ax, ay, Inches(1.7), Inches(0.52), RGBColor(0x40, 0x08, 0x15))
        add_text(s7, item, ax + Inches(0.05), ay + Inches(0.06), Inches(1.6), Inches(0.4),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        ax += Inches(1.78)
    else:
        add_text(s7, item, ax, ay + Inches(0.08), Inches(0.45), Inches(0.38),
                 font_size=16, bold=True, color=GOLD_LITE, align=PP_ALIGN.CENTER)
        ax += Inches(0.5)

slide_number(s7, 7, 10, light=True)

# ═══════════════════════════════════════════════════
#  SLIDE 8 — METRICS
# ═══════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, W, H, OFF_WHITE)

eyebrow(s8, "Performance & Portée", Inches(0.5), Inches(0.3), color=PRIMARY)
add_text(s8, "Chiffres Clés", Inches(0.5), Inches(0.6),
         Inches(7), Inches(0.7), font_size=32, bold=True, color=TEXT_DARK)

metrics = [
    ("5",      "",    "Écrans principaux\nentièrement fonctionnels"),
    ("20",     "+",   "DABs ATB & autres\nbanques pré-chargés"),
    ("3",      "",    "Langues : Derja,\nFrançais, Anglais"),
    ("0",      " €",  "Coût API maps\n& chatbot (gratuit)"),
]
mw = Inches(2.95)
mh = Inches(2.2)
mx0 = Inches(0.55)
my0 = Inches(1.5)
for i, (num, unit, lbl) in enumerate(metrics):
    mx2 = mx0 + i * (mw + Inches(0.2))
    add_rect(s8, mx2, my0, mw, mh, WHITE)
    add_rect(s8, mx2, my0, mw, Inches(0.06), PRIMARY)
    add_text(s8, num + unit, mx2 + Inches(0.2), my0 + Inches(0.25),
             mw - Inches(0.25), Inches(0.9), font_size=42, bold=True, color=PRIMARY)
    add_text(s8, lbl, mx2 + Inches(0.2), my0 + Inches(1.25),
             mw - Inches(0.25), Inches(0.8), font_size=12, color=DARK_GRAY, wrap=True)

# Progress bars
progs = [
    ("Fonctionnalités complètes", 95),
    ("Performance UI",            90),
    ("Couverture DABs Tunis",     80),
]
pw = Inches(3.85)
ph_start = Inches(4.0)
px3 = Inches(0.55)
for i, (lbl, pct) in enumerate(progs):
    px4 = px3 + i * (pw + Inches(0.2))
    add_rect(s8, px4, ph_start, pw, Inches(1.35), WHITE)
    add_text(s8, lbl, px4 + Inches(0.15), ph_start + Inches(0.12),
             pw - Inches(1.0), Inches(0.4), font_size=12, bold=False, color=TEXT_DARK)
    add_text(s8, f"{pct}%", px4 + pw - Inches(0.7), ph_start + Inches(0.12),
             Inches(0.6), Inches(0.4), font_size=13, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT)
    # bar bg
    add_rect(s8, px4 + Inches(0.15), ph_start + Inches(0.72), pw - Inches(0.28), Inches(0.14), GRAY_LIGHT)
    # bar fill
    add_rect(s8, px4 + Inches(0.15), ph_start + Inches(0.72),
             (pw - Inches(0.28)) * pct / 100, Inches(0.14), PRIMARY)

slide_number(s8, 8, 10, light=False)

# ═══════════════════════════════════════════════════
#  SLIDE 9 — ROADMAP
# ═══════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, W, H, WHITE)
add_rect(s9, 0, 0, W, Inches(0.06), GOLD)

eyebrow(s9, "Évolution du Projet", Inches(0.5), Inches(0.25), color=PRIMARY)
add_text(s9, "Roadmap de Développement", Inches(0.5), Inches(0.55),
         Inches(10), Inches(0.65), font_size=30, bold=True, color=TEXT_DARK)

phases = [
    ("✅", PRIMARY,      "Phase 1 · Complété",  "Infrastructure & Core UI",
     ["Navigation 5 onglets", "Thème ATB burgundy", "Modèles de données", "Android permissions"]),
    ("✅", PRIMARY,      "Phase 2 · Complété",  "Carte DABs & Assistant IA",
     ["Flutter Map + OSM", "Overpass API", "Chatbot Groq LLaMA", "Icône personnalisée"]),
    ("⏳", GOLD,         "Phase 3 · En cours",  "Optimisation & Polish",
     ["Tests utilisateurs", "Animations", "Biométrie", "Mode hors-ligne"]),
    ("🚀", DARK_GRAY,   "Phase 4 · Prévu",     "Intégration Backend ATB & Production",
     ["API bancaire réelle", "Paiements en ligne", "Play Store", "Analytics"]),
]

ry = Inches(1.45)
for (dot, dot_clr, date, title, items) in phases:
    # Connector line
    add_rect(s9, Inches(0.72), ry + Inches(0.48), Inches(0.05), Inches(1.4),
             GRAY_LIGHT if dot == "🚀" else PRIMARY)
    # Dot
    add_rect(s9, Inches(0.5), ry, Inches(0.48), Inches(0.48), dot_clr)
    add_text(s9, dot, Inches(0.5), ry, Inches(0.48), Inches(0.48),
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    # Content
    add_text(s9, date, Inches(1.2), ry - Inches(0.02), Inches(3.5), Inches(0.32),
             font_size=10, bold=True, color=PRIMARY)
    add_text(s9, title, Inches(1.2), ry + Inches(0.25), Inches(5), Inches(0.38),
             font_size=14, bold=True, color=TEXT_DARK)
    # Tags
    tx = Inches(6.8)
    for tag in items:
        tw2 = Inches(1.5)
        tag_bg = RGBColor(0xFE, 0xF0, 0xF2) if dot == "✅" else (
                 RGBColor(0xFF, 0xF9, 0xE6) if dot == "⏳" else GRAY_LIGHT)
        tag_c = PRIMARY if dot == "✅" else (RGBColor(0x7A, 0x5C, 0x00) if dot == "⏳" else DARK_GRAY)
        add_rect(s9, tx, ry + Inches(0.22), tw2, Inches(0.3), tag_bg)
        add_text(s9, tag, tx + Inches(0.08), ry + Inches(0.24), tw2 - Inches(0.1), Inches(0.26),
                 font_size=9, bold=True, color=tag_c)
        tx += tw2 + Inches(0.12)
    ry += Inches(1.5)

slide_number(s9, 9, 10, light=False)

# ═══════════════════════════════════════════════════
#  SLIDE 10 — THANK YOU
# ═══════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_rect(s10, 0, 0, W, H, PRIMARY_DARK)
add_rect(s10, 0, 0, W, H, PRIMARY)
add_rect(s10, 0, 0, W, Inches(0.06), GOLD)
add_rect(s10, 0, H - Inches(0.06), W, Inches(0.06), GOLD)

# ATB logo box
lw, lh = Inches(1.2), Inches(1.2)
lx2 = (W - lw) / 2
ly2 = Inches(0.9)
add_rect(s10, lx2, ly2, lw, lh, RGBColor(0x6A, 0x10, 0x22))
add_text(s10, "ATB", lx2, ly2 + Inches(0.28), lw, Inches(0.65),
         font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
add_textbox_ml(s10, [
    ("Merci.", 52, True, WHITE, PP_ALIGN.CENTER),
    ("Baraka Lahu Fik.", 52, True, GOLD_LITE, PP_ALIGN.CENTER),
], Inches(1), Inches(2.3), W - Inches(2), Inches(2.0))

# Subtitle
add_text(s10, "L'application ATB Mobile est prête — moderne, intelligente,\net pensée pour chaque client tunisien.",
         Inches(2), Inches(4.35), W - Inches(4), Inches(0.85),
         font_size=15, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# Buttons
add_rect(s10, Inches(3.5), Inches(5.4), Inches(2.5), Inches(0.55), GOLD)
add_text(s10, "▶  Voir la Démo", Inches(3.5), Inches(5.44), Inches(2.5), Inches(0.48),
         font_size=13, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

add_rect(s10, Inches(6.2), Inches(5.4), Inches(2.5), Inches(0.55), RGBColor(0x6A, 0x10, 0x22))
add_text(s10, "{ }  Code Source", Inches(6.2), Inches(5.44), Inches(2.5), Inches(0.48),
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Contact row
contacts = ["✉️  1hamza.achour@gmail.com", "🐦  Flutter 3.44", "📅  Juin 2026"]
cx3 = Inches(1.2)
for c in contacts:
    add_text(s10, c, cx3, Inches(6.3), Inches(3.5), Inches(0.4),
             font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    cx3 += Inches(3.6)

slide_number(s10, 10, 10, light=True)

# ── Save ────────────────────────────────────────────
out = r"C:\Users\MSI\AndroidStudioProjects\atb_banking_app\ATB_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
